#!/usr/bin/env python3

import os
import glob
import argparse
import pandas as pd
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# -----------------------------
# CLI
# -----------------------------
def parse_args():
    ap = argparse.ArgumentParser()
    ap.add_argument("--root", required=True, help="Root folder containing z<runno> subfolders")
    ap.add_argument("--excel", required=True, help="Excel file path")
    ap.add_argument("--out", required=True, help="Output PPTX filename")
    ap.add_argument("--template", default=None, help="Optional PPTX template to inherit theme/layout/size from")
    ap.add_argument("--slide_width", type=float, default=13.333, help='Slide width in inches (default: 13.333 for 16:9)')
    ap.add_argument("--slide_height", type=float, default=7.5, help='Slide height in inches (default: 7.5 for 16:9)')
    return ap.parse_args()


# -----------------------------
# Excel lookups
# -----------------------------
def load_lookup_table(excel_path: str) -> pd.DataFrame:
    df = pd.read_excel(excel_path)

    required = ["Arunno_or_Crunno", "Image_QA"]
    missing = [c for c in required if c not in df.columns]
    if missing:
        raise RuntimeError(f"Missing required columns in Excel: {missing}")

    df["Arunno_or_Crunno"] = df["Arunno_or_Crunno"].astype(str)
    df["Image_QA"] = df["Image_QA"].astype(str)
    return df


def lookup_row(df: pd.DataFrame, runno: str):
    matches = df[df["Arunno_or_Crunno"] == runno]
    if len(matches) == 0:
        return None
    return matches.iloc[0]


def lookup_status(df: pd.DataFrame, runno: str) -> str:
    row = lookup_row(df, runno)
    if row is None:
        return "Maybe"

    qa = str(row["Image_QA"]).strip().upper()
    if qa.startswith("U"):
        return "Yes"
    if qa.startswith("N"):
        return "No"
    return "Maybe"


def lookup_sex(df: pd.DataFrame, runno: str) -> str:
    """
    Returns 'm', 'f', or '?'.
    Looks for column 'Sex' (case-insensitive).
    Accepts M/F, Male/Female, etc.
    """
    row = lookup_row(df, runno)
    if row is None:
        return "?"

    sex_col = None
    for c in df.columns:
        if str(c).strip().lower() == "sex":
            sex_col = c
            break
    if sex_col is None:
        return "?"

    raw = str(row[sex_col]).strip().lower()
    if raw in ("m", "male"):
        return "m"
    if raw in ("f", "female"):
        return "f"
    return "?"


def parse_genotype_method(runno: str):
    prefix = runno[0].upper()
    if prefix == "A":
        return "APOE", "IV"
    if prefix == "P":
        return "APOE", "CM"
    if prefix == "C":
        return "CVN", "IV"
    if prefix == "V":
        return "CVN", "CM"
    return "Unknown", "Unknown"


# -----------------------------
# PPT helpers
# -----------------------------
def force_white_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)


def first_match(pattern: str):
    matches = sorted(glob.glob(pattern))
    return matches[0] if matches else None


def add_image(slide, pattern, left, top, width, height):
    path = first_match(pattern)
    if not path:
        print(f"Missing image: {pattern}")
        return False
    slide.shapes.add_picture(
        path,
        Inches(left),
        Inches(top),
        width=Inches(width),
        height=Inches(height),
    )
    return True


def add_textbox(
    slide,
    text,
    left,
    top,
    width,
    height,
    font_name,
    font_size,
    color_rgb,
    bold=False,
    align="left",      # "left" | "center" | "right"
    valign="top",      # "top" | "middle" | "bottom"
):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.clear()

    # vertical alignment
    if valign == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif valign == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    else:
        tf.vertical_anchor = MSO_ANCHOR.TOP

    p = tf.paragraphs[0]
    p.text = text

    # horizontal alignment
    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    elif align == "right":
        p.alignment = PP_ALIGN.RIGHT
    else:
        p.alignment = PP_ALIGN.LEFT

    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color_rgb

    return box	

# -----------------------------
# Title slide
# -----------------------------
def add_title_slide(prs, blank_layout):
    slide = prs.slides.add_slide(blank_layout)
    force_white_background(slide)

    add_textbox(
        slide,
        "Brain Clearance in Mice",
        left=0.9,
        top=1.6,
        width=11.6,
        height=1.0,
        font_name="Aptos Display",
        font_size=44,
        color_rgb=RGBColor(0, 0, 0),
        bold=True,
    )

    add_textbox(
        slide,
        "Initial Quality Assurance Report",
        left=0.95,
        top=2.6,
        width=11.6,
        height=0.6,
        font_name="Aptos Display",
        font_size=24,
        color_rgb=RGBColor(0, 0, 0),
        bold=False,
    )

    date_str = datetime.now().strftime("%B %d, %Y")
    add_textbox(
        slide,
        date_str,
        left=0.95,
        top=3.35,
        width=6.5,
        height=0.45,
        font_name="Aptos (Body)",
        font_size=16,
        color_rgb=RGBColor(0, 0, 0),
        bold=False,
    )

    add_textbox(
        slide,
        "Dr. B.J. Anderson, Ph.D., QIAL, Duke University Medical Center",
        left=0.95,
        top=4.0,
        width=12.0,
        height=0.6,
        font_name="Aptos (Body)",
        font_size=16,
        color_rgb=RGBColor(0, 0, 0),
        bold=False,
    )


# -----------------------------
# Summary counts + rubric helpers
# -----------------------------
def init_summary_counts():
    # counts[status][genotype][method] -> dict(sex->count, total->count)
    counts = {}
    for status in ("Yes", "Maybe", "No"):
        counts[status] = {}
        for genotype in ("APOE", "CVN"):
            counts[status][genotype] = {}
            for method in ("IV", "CM"):
                counts[status][genotype][method] = {"m": 0, "f": 0, "?": 0, "total": 0}
    return counts


def bump_counts(counts, status, genotype, method, sex):
    if status not in counts:
        return
    if genotype not in counts[status]:
        return
    if method not in counts[status][genotype]:
        return
    if sex not in ("m", "f", "?"):
        sex = "?"
    cell = counts[status][genotype][method]
    cell["total"] += 1
    cell[sex] += 1

def add_cell(slide, left, top, width, height, border_width_pt=1.5):
    """
    Draw a border-only rectangle (no fill) to act as a table cell boundary.
    """
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    # No fill
    shape.fill.background()

    # Border
    shape.line.color.rgb = RGBColor(0, 0, 0)
    shape.line.width = Pt(border_width_pt)
    return shape


def fmt_cell(cell: dict) -> str:
    total = cell.get("total", 0)
    m = cell.get("m", 0)
    f = cell.get("f", 0)
    q = cell.get("?", 0)
    bits = [f"{m} m", f"{f} f"]
    if q:
        bits.append(f"{q} ?")
    return f"{total} ({', '.join(bits)})"


def add_rubric(slide, title, counts_block, left, top, width, height, big=False):
    """
    Renders a 2x2 rubric with borders:
      columns: APOE, CVN
      rows: IV, CM
    """

    # Typography: labels >= counts
    count_fs = 24 if big else 16
    label_fs = count_fs + (2 if big else 2)  # labels slightly bigger
    title_fs = 28 if big else 20
    title_h = 0.5 if big else 0.4

    # Title
    add_textbox(
        slide,
        title,
        left=left,
        top=top,
        width=width,
        height=title_h,
        font_name="Aptos Display",
        font_size=title_fs,
        color_rgb=RGBColor(0, 0, 0),
        bold=True,
        align="left",
        valign="middle",
    )

    # Table geometry
    grid_top = top + (0.65 if big else 0.55)
    grid_h = height - (0.85 if big else 0.7)

    label_col_w = 1.0 if big else 0.85
    data_col_w = (width - label_col_w) / 2.0

    header_h = 0.65 if big else 0.5
    row_h = (grid_h - header_h) / 2.0

    # --- Draw borders (header row + 2 data rows, 3 columns) ---
    # Row 0: header
    add_cell(slide, left, grid_top, label_col_w, header_h, border_width_pt=2 if big else 1.5)
    add_cell(slide, left + label_col_w, grid_top, data_col_w, header_h, border_width_pt=2 if big else 1.5)
    add_cell(slide, left + label_col_w + data_col_w, grid_top, data_col_w, header_h, border_width_pt=2 if big else 1.5)

    # Row 1: IV
    add_cell(slide, left, grid_top + header_h, label_col_w, row_h, border_width_pt=2 if big else 1.5)
    add_cell(slide, left + label_col_w, grid_top + header_h, data_col_w, row_h, border_width_pt=2 if big else 1.5)
    add_cell(slide, left + label_col_w + data_col_w, grid_top + header_h, data_col_w, row_h, border_width_pt=2 if big else 1.5)

    # Row 2: CM
    add_cell(slide, left, grid_top + header_h + row_h, label_col_w, row_h, border_width_pt=2 if big else 1.5)
    add_cell(slide, left + label_col_w, grid_top + header_h + row_h, data_col_w, row_h, border_width_pt=2 if big else 1.5)
    add_cell(slide, left + label_col_w + data_col_w, grid_top + header_h + row_h, data_col_w, row_h, border_width_pt=2 if big else 1.5)

    # --- Header labels ---
    add_textbox(
        slide, "", left, grid_top, label_col_w, header_h,
        "Aptos Display", label_fs, RGBColor(0, 0, 0),
        bold=True, align="center", valign="middle"
    )
    add_textbox(
        slide, "APOE", left + label_col_w, grid_top, data_col_w, header_h,
        "Aptos Display", label_fs, RGBColor(0, 0, 0),
        bold=True, align="center", valign="middle"
    )
    add_textbox(
        slide, "CVN", left + label_col_w + data_col_w, grid_top, data_col_w, header_h,
        "Aptos Display", label_fs, RGBColor(0, 0, 0),
        bold=True, align="center", valign="middle"
    )

    # --- Row labels ---
    add_textbox(
        slide, "IV", left, grid_top + header_h, label_col_w, row_h,
        "Aptos Display", label_fs, RGBColor(0, 0, 0),
        bold=True, align="center", valign="middle"
    )
    add_textbox(
        slide, "CM", left, grid_top + header_h + row_h, label_col_w, row_h,
        "Aptos Display", label_fs, RGBColor(0, 0, 0),
        bold=True, align="center", valign="middle"
    )

    # --- Data cells ---
    add_textbox(
        slide,
        fmt_cell(counts_block["APOE"]["IV"]),
        left + label_col_w, grid_top + header_h, data_col_w, row_h,
        "Aptos (Body)", count_fs, RGBColor(0, 0, 0),
        bold=False, align="center", valign="middle"
    )
    add_textbox(
        slide,
        fmt_cell(counts_block["CVN"]["IV"]),
        left + label_col_w + data_col_w, grid_top + header_h, data_col_w, row_h,
        "Aptos (Body)", count_fs, RGBColor(0, 0, 0),
        bold=False, align="center", valign="middle"
    )
    add_textbox(
        slide,
        fmt_cell(counts_block["APOE"]["CM"]),
        left + label_col_w, grid_top + header_h + row_h, data_col_w, row_h,
        "Aptos (Body)", count_fs, RGBColor(0, 0, 0),
        bold=False, align="center", valign="middle"
    )
    add_textbox(
        slide,
        fmt_cell(counts_block["CVN"]["CM"]),
        left + label_col_w + data_col_w, grid_top + header_h + row_h, data_col_w, row_h,
        "Aptos (Body)", count_fs, RGBColor(0, 0, 0),
        bold=False, align="center", valign="middle"
    )
    # Title for the block
    add_textbox(
        slide,
        title,
        left=left,
        top=top,
        width=width,
        height=0.45 if big else 0.38,
        font_name="Aptos Display",
        font_size=26 if big else 18,
        color_rgb=RGBColor(0, 0, 0),
        bold=True,
    )

    grid_top = top + (0.6 if big else 0.5)
    grid_h = height - (0.8 if big else 0.65)

    # Grid geometry: left label col + 2 data cols; header row + 2 data rows
    label_col_w = 0.9 if big else 0.75
    data_col_w = (width - label_col_w) / 2.0

    header_h = 0.55 if big else 0.45
    row_h = (grid_h - header_h) / 2.0

    # Column headers
    add_textbox(slide, "", left, grid_top, label_col_w, header_h, "Aptos (Body)", 14 if big else 12, RGBColor(0, 0, 0), False)
    add_textbox(slide, "APOE", left + label_col_w, grid_top, data_col_w, header_h, "Aptos Display", 16 if big else 14, RGBColor(0, 0, 0), True)
    add_textbox(slide, "CVN", left + label_col_w + data_col_w, grid_top, data_col_w, header_h, "Aptos Display", 16 if big else 14, RGBColor(0, 0, 0), True)

    # Row labels
    add_textbox(slide, "IV", left, grid_top + header_h, label_col_w, row_h, "Aptos Display", 16 if big else 14, RGBColor(0, 0, 0), True)
    add_textbox(slide, "CM", left, grid_top + header_h + row_h, label_col_w, row_h, "Aptos Display", 16 if big else 14, RGBColor(0, 0, 0), True)

    # Data cells
    # IV row
    add_textbox(
        slide,
        fmt_cell(counts_block["APOE"]["IV"]),
        left + label_col_w,
        grid_top + header_h,
        data_col_w,
        row_h,
        "Aptos (Body)",
        22 if big else 14,
        RGBColor(0, 0, 0),
        False,
    )
    add_textbox(
        slide,
        fmt_cell(counts_block["CVN"]["IV"]),
        left + label_col_w + data_col_w,
        grid_top + header_h,
        data_col_w,
        row_h,
        "Aptos (Body)",
        22 if big else 14,
        RGBColor(0, 0, 0),
        False,
    )

    # CM row
    add_textbox(
        slide,
        fmt_cell(counts_block["APOE"]["CM"]),
        left + label_col_w,
        grid_top + header_h + row_h,
        data_col_w,
        row_h,
        "Aptos (Body)",
        22 if big else 14,
        RGBColor(0, 0, 0),
        False,
    )
    add_textbox(
        slide,
        fmt_cell(counts_block["CVN"]["CM"]),
        left + label_col_w + data_col_w,
        grid_top + header_h + row_h,
        data_col_w,
        row_h,
        "Aptos (Body)",
        22 if big else 14,
        RGBColor(0, 0, 0),
        False,
    )


def add_summary_slide(prs, blank_layout, summary_counts):
    slide = prs.slides.add_slide(blank_layout)
    force_white_background(slide)

    # Slide title
    add_textbox(
        slide,
        "QA Summary",
        left=0.9,
        top=0.5,
        width=11.6,
        height=0.7,
        font_name="Aptos Display",
        font_size=40,
        color_rgb=RGBColor(0, 0, 0),
        bold=True,
    )

    # Layout:
    # Left half (big): Yes
    # Right half (stacked): Maybe (top), No (bottom)
    left_x = 0.9
    left_w = 6.3
    right_x = 7.4
    right_w = 5.5

    # Big left rubric (Usable)
    add_rubric(
        slide,
        "Usable Data",
        summary_counts["Yes"],
        left=left_x,
        top=1.45,
        width=left_w,
        height=5.8,
        big=True,
    )

    # Right top rubric (Maybe)
    add_rubric(
        slide,
        "Possibly Usable Data",
        summary_counts["Maybe"],
        left=right_x,
        top=1.45,
        width=right_w,
        height=2.8,
        big=False,
    )

    # Right bottom rubric (No)
    add_rubric(
        slide,
        "Unusable Data",
        summary_counts["No"],
        left=right_x,
        top=4.35,
        width=right_w,
        height=2.9,
        big=False,
    )

    # Footnote
    add_textbox(
        slide,
        "*Data may be improved/made usable after motion corrections (coregistration and volume pruning).",
        left=0.9,
        top=7.05,
        width=12.4,
        height=0.35,
        font_name="Aptos (Body)",
        font_size=12,
        color_rgb=RGBColor(0, 0, 0),
        bold=False,
    )


# -----------------------------
# Main
# -----------------------------
def main():
    args = parse_args()
    df = load_lookup_table(args.excel)

    # Pass 1: scan folders, collect run list + summary counts
    run_folders = []
    summary_counts = init_summary_counts()

    for folder in sorted(os.listdir(args.root)):
        if not folder.startswith("z"):
            continue
        runno = folder[1:]
        fullpath = os.path.join(args.root, folder)
        if not os.path.isdir(fullpath):
            continue
        if not glob.glob(os.path.join(fullpath, "*.png")):
            continue

        status = lookup_status(df, runno)
        sex = lookup_sex(df, runno)
        genotype, method = parse_genotype_method(runno)

        bump_counts(summary_counts, status, genotype, method, sex)
        run_folders.append((runno, fullpath))

    # Pass 2: build PPT in the required order
    if args.template:
        prs = Presentation(args.template)
        # template controls size/theme
    else:
        prs = Presentation()
        prs.slide_width = Inches(args.slide_width)
        prs.slide_height = Inches(args.slide_height)

    blank_layout = prs.slide_layouts[6]

    # Title then Summary immediately after
    add_title_slide(prs, blank_layout)
    add_summary_slide(prs, blank_layout, summary_counts)

    # Then the per-run slides
    for runno, fullpath in run_folders:
        genotype, method = parse_genotype_method(runno)
        status = lookup_status(df, runno)

        print(f"Adding slide for {runno}")

        slide = prs.slides.add_slide(blank_layout)
        force_white_background(slide)

        # -------- Images --------
        add_image(slide, os.path.join(fullpath, f"{runno}_CM_roi_intensity_xyz_*_r2.5.png"), 0.1, 0.0, 4.25, 3.19)
        add_image(slide, os.path.join(fullpath, f"{runno}_CSF_roi_intensity_xyz_*_r2.5.png"), 4.55, 0.0, 4.25, 3.19)
        add_image(slide, os.path.join(fullpath, f"{runno}_Hc_roi_intensity_xyz_*_r2.5.png"), 9.0, 0.0, 4.25, 3.19)

        add_image(slide, os.path.join(fullpath, f"{runno}_CM.png"), 0.1, 3.16, 4.25, 4.25)
        add_image(slide, os.path.join(fullpath, f"{runno}_CSF.png"), 4.55, 3.16, 4.25, 4.25)
        add_image(slide, os.path.join(fullpath, f"{runno}_Hc.png"), 9.0, 3.16, 4.25, 4.25)

        # -------- Bottom Labels (Set 1) --------
        add_textbox(slide, "Cisterna Magna", 2.6, 7.0, 1.41, 0.3, "Aptos (Body)", 12, RGBColor(255, 255, 255))
        add_textbox(slide, "Cerebral Spinal Fluid", 6.95, 7.0, 1.97, 0.3, "Aptos (Body)", 12, RGBColor(255, 255, 255))
        add_textbox(slide, "Cisterna Magna", 11.7, 7.0, 1.41, 0.3, "Aptos (Body)", 12, RGBColor(255, 255, 255))

        # -------- Genotype/Method/Usable (Set 2) [YOUR UPDATED POSITIONS] --------
        add_textbox(slide, f"Genotype: {genotype}", 3.5, 2.85, 1.83, 0.37, "Aptos Display", 16, RGBColor(0, 0, 0))
        add_textbox(slide, f"Method: {method}", 8.1, 2.85, 1.83, 0.37, "Aptos Display", 16, RGBColor(0, 0, 0))
        add_textbox(slide, f"Usable: {status}", 11.7, 2.85, 1.83, 0.37, "Aptos Display", 16, RGBColor(0, 0, 0))

        # -------- Title box ($runno) --------
        add_textbox(slide, runno, 0.1, 2.85, 1.5, 0.4, "Aptos Display", 16, RGBColor(0, 0, 0), bold=True)

    prs.save(args.out)
    print(f"\nSaved presentation to {args.out}")


if __name__ == "__main__":
    main()